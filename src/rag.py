import io
import mimetypes
import json
from typing import List
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
import openpyxl
import chromadb
from chromadb.utils import embedding_functions

# -----------------------
# Config
# -----------------------
CHUNK_SIZE = 600
CHUNK_OVERLAP = 150
EMBED_MODEL = "all-MiniLM-L6-v2"

# -----------------------
# Model & Vector Store
# -----------------------
client = chromadb.PersistentClient(path="chroma_data")  # kalıcı depolama
sentence_transformer_ef = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name=EMBED_MODEL
)
collection = client.get_or_create_collection(
    name="knowledge_bot",
    embedding_function=sentence_transformer_ef
)

# -----------------------
# Helpers
# -----------------------
def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Metni belirli uzunluklarda chunklara ayırır."""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        chunks.append(text[start:end])
        start += size - overlap
    return chunks

# -----------------------
# Dosya Parsing
# -----------------------
def extract_text_from_file(filename: str, raw: bytes) -> str:
    """PDF, Word, PPTX, Excel, JSON, TXT vb. dosyalardan metin çıkarır."""
    text = ""
    mime_type, _ = mimetypes.guess_type(filename)

    try:
        if mime_type == "application/pdf":
            try:
                pdf = PdfReader(io.BytesIO(raw))
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            except Exception:
                images = convert_from_bytes(raw)
                for img in images:
                    text += pytesseract.image_to_string(img, lang="tur") + "\n"

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(io.BytesIO(raw))
            for p in doc.paragraphs:
                text += p.text + "\n"
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_data = rel.target_part.blob
                    img = Image.open(io.BytesIO(image_data))
                    text += pytesseract.image_to_string(img, lang="tur") + "\n"

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(raw))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    if shape.shape_type == 13:  # Picture
                        image = Image.open(io.BytesIO(shape.image.blob))
                        text += pytesseract.image_to_string(image, lang="tur") + "\n"

        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or filename.endswith(".xls"):
            wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
            for sheet in wb.worksheets:
                text += f"\n# {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + "\n"            

        elif filename.endswith(".json"):
            data = json.loads(raw.decode("utf-8", errors="ignore"))
            if isinstance(data, dict) and "items" in data:
                for item in data["items"]:
                    if "q" in item:
                        text += item["q"] + "\n"
                    if "a" in item:
                        text += item["a"] + "\n"
            else:
                text = json.dumps(data)

        elif filename.endswith(".md") or filename.endswith(".txt"):
            text = raw.decode("utf-8", errors="ignore")

        else:
            text = raw.decode("utf-8", errors="ignore")

    except Exception as e:
        text = f"[extract_text_from_file error: {str(e)}]"

    return text.strip()

# -----------------------
# Indexleme
# -----------------------
def index_doc(filename: str, text: str, topic: str = "other") -> int:
    """Metni chunklara bölerek Chroma koleksiyonuna ekler."""
    chunks = chunk_text(text)
    for i, chunk in enumerate(chunks):
        collection.add(
            ids=[f"{filename}_{i}"],
            documents=[chunk],
            metadatas=[{"doc_id": filename, "chunk": i, "topic": topic}]
        )
    return len(chunks)

# -----------------------
# Arama
# -----------------------
def search(query: str, top_k: int = 5, where: dict = None):
    """Sorgu ile Chroma koleksiyonunda arama yapar."""
    results = collection.query(
        query_texts=[query],
        n_results=top_k,
        where=where
    )
    return results

# -----------------------
# Silme
# -----------------------
def delete_doc(doc_id: str):
    try:
        collection.delete(where={"doc_id": doc_id})
        return {"status": "deleted", "doc_id": doc_id}
    except Exception as e:
        return {"status": "error", "detail": str(e)}

def delete_all():
    try:
        client.delete_collection("knowledge_bot")
        global collection
        collection = client.get_or_create_collection(
            name="knowledge_bot",
            embedding_function=sentence_transformer_ef
        )
        return {"status": "all deleted"}
    except Exception as e:
        return {"status": "error", "detail": str(e)}
